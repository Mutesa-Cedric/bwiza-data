"""Tests for Office document extraction (DOCX, PPTX, DOC)."""

from __future__ import annotations

import io
from unittest.mock import patch

from apps.common.office import (
    _extract_images_from_zip,
    extract_office_text,
    is_office_type,
    ocr_office,
)

# ── is_office_type ──────────────────────────────────────────


def test_is_office_type_docx():
    assert is_office_type(
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )


def test_is_office_type_pptx():
    assert is_office_type(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


def test_is_office_type_doc():
    assert is_office_type("application/msword")


def test_is_office_type_rejects_pdf():
    assert not is_office_type("application/pdf")


def test_is_office_type_rejects_html():
    assert not is_office_type("text/html")


# ── extract_office_text ─────────────────────────────────────


def _make_docx(text: str) -> bytes:
    """Create a minimal DOCX file with the given text."""
    from docx import Document

    doc = Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(texts: list[str]) -> bytes:
    """Create a minimal PPTX file with one slide per text."""
    from pptx import Presentation

    prs = Presentation()
    for text in texts:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.placeholders[1].text = text  # type: ignore[attr-defined]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_docx():
    text = "Mu Rwanda uburezi ni ingenzi cyane ku iterambere ry'igihugu. " * 5
    data = _make_docx(text)
    ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = extract_office_text(data, ct, url="test.docx")
    assert result is not None
    assert "uburezi" in result.text


def test_extract_pptx():
    texts = [
        "Mu Rwanda uburezi ni ingenzi cyane ku iterambere ry'igihugu. " * 3,
        "Abanyeshuri biga amasomo atandukanye harimo ikinyarwanda. " * 3,
    ]
    data = _make_pptx(texts)
    ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    result = extract_office_text(data, ct, url="test.pptx")
    assert result is not None
    assert "uburezi" in result.text


def test_extract_office_too_short():
    data = _make_docx("Hi")
    ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = extract_office_text(data, ct, url="short.docx")
    assert result is None


def test_extract_office_unknown_type():
    result = extract_office_text(b"data", "application/octet-stream", url="file.bin")
    assert result is None


def test_extract_office_corrupt_file():
    ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = extract_office_text(b"not-a-real-docx", ct, url="corrupt.docx")
    assert result is None


# ── _extract_images_from_zip ────────────────────────────────


def test_extract_images_from_valid_zip():
    import zipfile

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("word/media/image1.png", b"\x89PNG fake image data")
        zf.writestr("word/media/image2.jpg", b"\xff\xd8 fake jpeg data")
        zf.writestr("word/document.xml", b"<xml>not an image</xml>")
    images = _extract_images_from_zip(buf.getvalue())
    assert len(images) == 2


def test_extract_images_from_bad_zip():
    images = _extract_images_from_zip(b"not a zip file")
    assert images == []


# ── ocr_office ──────────────────────────────────────────────


@patch("apps.common.ocr._OCR_BACKEND", "none")
def test_ocr_office_no_backend():
    ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = ocr_office(b"data", ct)
    assert result is None


def test_ocr_office_unsupported_format():
    result = ocr_office(b"data", "application/msword")
    assert result is None


@patch("apps.common.ocr._OCR_BACKEND", "apple_vision")
def test_ocr_office_no_images():
    """DOCX with no embedded images returns None."""
    # Make a real DOCX with text only (no images)
    data = _make_docx("Just text " * 20)
    ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = ocr_office(data, ct)
    assert result is None
