"""Extract text from Office documents (DOCX, PPTX, DOC).

Falls back to OCR for scanned/image-only content embedded in documents.
"""

from __future__ import annotations

import io
import zipfile

from apps.common.logging import get_logger
from apps.targeted_crawler.extract import ExtractedDoc

log = get_logger(__name__)

# Image extensions we can OCR from embedded office content
_IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp")

# Content type → format mapping
_OFFICE_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/msword": "doc",
}

# Minimum chars to consider extraction successful
_MIN_TEXT_CHARS = 50


def is_office_type(content_type: str) -> bool:
    """Check if content type is an Office document we can extract."""
    return any(ct in content_type for ct in _OFFICE_TYPES)


def _extract_docx(data: bytes) -> str:
    """Extract text from a DOCX file."""
    from docx import Document

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_pptx(data: bytes) -> str:
    """Extract text from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    texts: list[str] = []
    for slide in prs.slides:
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


def _extract_doc(data: bytes) -> str | None:
    """Extract text from a legacy DOC file using antiword or textract."""
    import subprocess
    import tempfile

    with tempfile.NamedTemporaryFile(suffix=".doc", delete=True) as tmp:
        tmp.write(data)
        tmp.flush()
        try:
            result = subprocess.run(
                ["antiword", tmp.name],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except FileNotFoundError:
            log.debug("antiword not installed, cannot extract .doc files")
        except subprocess.TimeoutExpired:
            log.debug("antiword timed out on .doc file")
    return None


def extract_office_text(
    data: bytes,
    content_type: str,
    url: str = "",
) -> ExtractedDoc | None:
    """Extract text from an Office document.

    Supports DOCX, PPTX, and DOC (via antiword). If text extraction
    yields too little content, returns None so the caller can try OCR.

    Returns ExtractedDoc or None if extraction fails.
    """
    fmt = None
    for ct, f in _OFFICE_TYPES.items():
        if ct in content_type:
            fmt = f
            break

    if fmt is None:
        return None

    try:
        if fmt == "docx":
            text = _extract_docx(data)
        elif fmt == "pptx":
            text = _extract_pptx(data)
        elif fmt == "doc":
            text = _extract_doc(data)
            if text is None:
                return None
        else:
            return None
    except Exception:
        log.debug("Office extraction failed for %s (%s)", url, fmt)
        return None

    if not text or len(text) < _MIN_TEXT_CHARS:
        log.debug("Office extraction: too little text (%d chars) from %s", len(text or ""), url)
        return None

    log.info("Office[%s]: extracted %d chars from %s", fmt, len(text), url)
    return ExtractedDoc(title="", text=text)


def _extract_images_from_zip(data: bytes) -> list[bytes]:
    """Extract image files from a ZIP-based Office document (DOCX/PPTX)."""
    images: list[bytes] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            for name in sorted(zf.namelist()):
                if name.lower().endswith(_IMAGE_EXTS):
                    images.append(zf.read(name))
    except (zipfile.BadZipFile, Exception):
        log.debug("Failed to extract images from office document")
    return images


def ocr_office(
    data: bytes,
    content_type: str,
    url: str = "",
    max_images: int = 200,
) -> ExtractedDoc | None:
    """OCR fallback for scanned Office documents (DOCX/PPTX).

    Extracts embedded images from the ZIP-based Office format and runs
    them through the OCR pipeline. Only works for DOCX/PPTX (not legacy DOC).

    Returns None if no images found or OCR yields insufficient text.
    """
    import apps.common.ocr as _ocr_mod

    if _ocr_mod._OCR_BACKEND == "none":
        log.debug("OCR: no backend available for office fallback")
        return None

    fmt = None
    for ct, f in _OFFICE_TYPES.items():
        if ct in content_type:
            fmt = f
            break

    if fmt not in ("docx", "pptx"):
        log.debug("OCR: office fallback not supported for %s", fmt or "unknown")
        return None

    images = _extract_images_from_zip(data)
    if not images:
        log.debug("OCR: no images found in %s (%s)", url, fmt)
        return None

    if len(images) > max_images:
        log.debug("OCR: too many images (%d) in %s, capping at %d", len(images), url, max_images)
        images = images[:max_images]

    from PIL import Image

    page_texts: list[str] = []
    for i, img_bytes in enumerate(images):
        try:
            img = Image.open(io.BytesIO(img_bytes))
            # Skip very small images (icons, logos)
            if img.width < 100 or img.height < 100:
                continue
            if _ocr_mod._OCR_BACKEND == "apple_vision":
                text = _ocr_mod._ocr_page_apple_vision(img)
            else:
                text = _ocr_mod._ocr_page_tesseract(img, "eng+fra")
        except Exception:
            log.debug("OCR: failed on image %d of %s", i + 1, url)
            continue

        if len(text) >= 30:
            page_texts.append(text)

    if not page_texts:
        log.debug("OCR: no text from %d images in %s", len(images), url)
        return None

    full_text = "\n\n".join(page_texts)
    log.info(
        "OCR[office/%s]: extracted %d chars from %d images: %s",
        fmt,
        len(full_text),
        len(page_texts),
        url,
    )
    return ExtractedDoc(title="", text=full_text)
